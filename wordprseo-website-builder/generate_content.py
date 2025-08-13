import os
import re
import sys


def read_file(path: str) -> str:
    """Extract text from various document types."""
    _, ext = os.path.splitext(path)
    ext = ext.lower()
    if ext == ".txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    if ext in {".docx", ".doc"}:
        try:
            from docx import Document

            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception:
            return ""
    if ext == ".pdf":
        try:
            import PyPDF2

            with open(path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return ""
    if ext in {".ppt", ".pptx"}:
        try:
            from pptx import Presentation

            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception:
            return ""
    return ""


def _sentences(text: str) -> list[str]:
    parts = re.split(r"(?<=[.!?])\s+", text.strip())
    return [p.strip() for p in parts if p.strip()]


def _chunk(sentences: list[str], start: int, size: int = 3) -> str:
    chunk = sentences[start : start + size]
    if len(chunk) < size:
        chunk += ["Placeholder content."] * (size - len(chunk))
    return " ".join(chunk)


def generate_homepage(text: str) -> str:
    """Generate homepage sections with titles and subtitles."""
    sents = _sentences(text)
    hero = _chunk(sents, 0)
    about = _chunk(sents, 3)
    services = _chunk(sents, 6)
    contact = _chunk(sents, 9)
    return (
        "# Hero\n"
        "## Welcome to Our Site\n"
        f"{hero}\n\n"
        "# About Us\n"
        "## Who We Are\n"
        f"{about}\n\n"
        "# Services\n"
        "## What We Do\n"
        f"{services}\n\n"
        "# Contact\n"
        "## Get in Touch\n"
        f"{contact}"
    )


def main() -> None:
    if len(sys.argv) < 3:
        print("Usage: generate_content.py <page> <file>", file=sys.stderr)
        sys.exit(1)
    page, file_path = sys.argv[1], sys.argv[2]
    text = read_file(file_path)
    if page == "home":
        print(generate_homepage(text))
    else:
        print(f"Page type '{page}' not supported.")


if __name__ == "__main__":
    main()
